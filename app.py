import io
import os
import json
from datetime import datetime
from functools import wraps

from flask import (
    Flask, request, jsonify, render_template,
    redirect, url_for, session, abort, send_file
)
from flask_sqlalchemy import SQLAlchemy

try:
    from anthropic import Anthropic
except ImportError:
    Anthropic = None

DB_PATH = os.environ.get("DB_PATH", os.path.join(os.path.dirname(__file__), "tempform.db"))

app = Flask(__name__)
app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", "tempform-dev-secret-change-me")
app.config["SQLALCHEMY_DATABASE_URI"] = f"sqlite:///{DB_PATH}"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False

db = SQLAlchemy(app)

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
anthropic_client = Anthropic(api_key=ANTHROPIC_API_KEY) if (Anthropic and ANTHROPIC_API_KEY) else None

MEMBERS = {
    "001": {"name": "劉浩晨", "role": "admin"},
    "002": {"name": "蔡崇正", "role": "member"},
    "003": {"name": "蕭淳勻", "role": "member"},
    "004": {"name": "洪瑞鶯", "role": "member"},
    "005": {"name": "蔡幸慧", "role": "member"},
}

GROUPS = [f"第{i}組" for i in range(1, 7)]
RUBRIC = ["內容深度", "實務連結", "批判反思", "整合脈絡", "表達呈現"]


class Submission(db.Model):
    __tablename__ = "submissions"
    id = db.Column(db.Integer, primary_key=True)
    member_code = db.Column(db.String(8), unique=True, nullable=False, index=True)
    member_name = db.Column(db.String(32), nullable=False)
    payload = db.Column(db.Text, nullable=False, default="{}")
    status = db.Column(db.String(16), nullable=False, default="draft")  # draft / submitted
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    def to_dict(self):
        return {
            "member_code": self.member_code,
            "member_name": self.member_name,
            "status": self.status,
            "payload": json.loads(self.payload or "{}"),
            "updated_at": self.updated_at.strftime("%Y-%m-%d %H:%M:%S") if self.updated_at else None,
        }


class AISummary(db.Model):
    __tablename__ = "ai_summaries"
    id = db.Column(db.Integer, primary_key=True)
    payload = db.Column(db.Text, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


with app.app_context():
    os.makedirs(os.path.dirname(DB_PATH) or ".", exist_ok=True)
    db.create_all()


def login_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if "code" not in session:
            return redirect(url_for("login"))
        return f(*args, **kwargs)
    return wrapper


def admin_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        code = session.get("code")
        if not code or MEMBERS.get(code, {}).get("role") != "admin":
            return abort(403)
        return f(*args, **kwargs)
    return wrapper


@app.route("/")
def index():
    code = session.get("code")
    if not code:
        return redirect(url_for("login"))
    if MEMBERS.get(code, {}).get("role") == "admin":
        return redirect(url_for("admin_page"))
    return redirect(url_for("form_page"))


def _is_admin():
    code = session.get("code")
    return bool(code) and MEMBERS.get(code, {}).get("role") == "admin"


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        code = (request.form.get("code") or "").strip()
        if code in MEMBERS:
            session["code"] = code
            session["name"] = MEMBERS[code]["name"]
            return redirect(url_for("index"))
        return render_template("login.html", members=MEMBERS, error="代號錯誤")
    return render_template("login.html", members=MEMBERS, error=None)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


@app.route("/form")
@login_required
def form_page():
    me = session["code"]
    target = request.args.get("as", "").strip()
    if target and target != me:
        if not _is_admin():
            return abort(403)
        if target not in MEMBERS:
            return abort(404)
        edit_code = target
        admin_editing = True
    else:
        edit_code = me
        admin_editing = False

    sub = Submission.query.filter_by(member_code=edit_code).first()
    payload = json.loads(sub.payload) if sub else {}
    status = sub.status if sub else "new"
    return render_template(
        "form.html",
        code=edit_code,
        name=MEMBERS[edit_code]["name"],
        groups=GROUPS,
        rubric=RUBRIC,
        payload=payload,
        status=status,
        admin_editing=admin_editing,
        viewer_is_admin=_is_admin(),
    )


def _save_submission(code, payload, status):
    sub = Submission.query.filter_by(member_code=code).first()
    if not sub:
        sub = Submission(member_code=code, member_name=MEMBERS[code]["name"])
        db.session.add(sub)
    sub.payload = json.dumps(payload, ensure_ascii=False)
    sub.member_name = MEMBERS[code]["name"]
    sub.status = status
    sub.updated_at = datetime.utcnow()
    db.session.commit()
    return sub


@app.route("/api/save", methods=["POST"])
@login_required
def api_save():
    me = session["code"]
    data = request.get_json(silent=True) or {}
    payload = data.get("payload") or {}
    action = data.get("action", "draft")  # draft / submit
    target = (data.get("target_code") or me).strip()

    if target != me:
        if not _is_admin():
            return jsonify({"error": "forbidden"}), 403
        if target not in MEMBERS:
            return jsonify({"error": "not found"}), 404

    status = "submitted" if action == "submit" else "draft"
    sub = _save_submission(target, payload, status)
    return jsonify({"ok": True, "status": status, "updated_at": sub.updated_at.strftime("%Y-%m-%d %H:%M:%S")})


def _empty_payload():
    return {
        "groups": {g: {"q1": "", "q2": "", "q3": {r: 0 for r in RUBRIC}} for g in GROUPS},
        "overall": {"best_group": "", "best_reason": "", "worst_group": "", "worst_reason": "", "self_review": ""},
    }


@app.route("/api/admin/member/<code>", methods=["DELETE"])
@admin_required
def api_admin_delete_member(code):
    if code not in MEMBERS:
        return jsonify({"error": "not found"}), 404
    sub = Submission.query.filter_by(member_code=code).first()
    if sub:
        db.session.delete(sub)
        db.session.commit()
    return jsonify({"ok": True})


@app.route("/api/admin/member/<code>/clear", methods=["POST"])
@admin_required
def api_admin_clear(code):
    if code not in MEMBERS:
        return jsonify({"error": "not found"}), 404
    data = request.get_json(silent=True) or {}
    scope = data.get("scope", "")  # "group:第1組" / "overall" / "all"
    sub = Submission.query.filter_by(member_code=code).first()
    if not sub:
        return jsonify({"ok": True, "noop": True})
    payload = json.loads(sub.payload or "{}")
    base = _empty_payload()
    payload.setdefault("groups", base["groups"])
    payload.setdefault("overall", base["overall"])

    if scope == "all":
        payload = base
    elif scope == "overall":
        payload["overall"] = base["overall"]
    elif scope.startswith("group:"):
        g = scope.split(":", 1)[1]
        if g not in GROUPS:
            return jsonify({"error": "bad group"}), 400
        payload["groups"][g] = {"q1": "", "q2": "", "q3": {r: 0 for r in RUBRIC}}
    else:
        return jsonify({"error": "bad scope"}), 400

    _save_submission(code, payload, sub.status)
    return jsonify({"ok": True})


@app.route("/admin")
@admin_required
def admin_page():
    return render_template("admin.html", groups=GROUPS, rubric=RUBRIC, members=MEMBERS)


@app.route("/api/admin/overview")
@admin_required
def api_admin_overview():
    rows = []
    subs = {s.member_code: s for s in Submission.query.all()}
    for code, info in MEMBERS.items():
        sub = subs.get(code)
        rows.append({
            "code": code,
            "name": info["name"],
            "role": info["role"],
            "status": sub.status if sub else "none",
            "updated_at": sub.updated_at.strftime("%Y-%m-%d %H:%M:%S") if sub and sub.updated_at else None,
        })
    return jsonify({"rows": rows})


@app.route("/api/admin/member/<code>")
@admin_required
def api_admin_member(code):
    if code not in MEMBERS:
        return jsonify({"error": "not found"}), 404
    sub = Submission.query.filter_by(member_code=code).first()
    if not sub:
        return jsonify({
            "code": code,
            "name": MEMBERS[code]["name"],
            "status": "none",
            "payload": {},
            "updated_at": None,
        })
    return jsonify(sub.to_dict())


def _avg(nums):
    nums = [n for n in nums if isinstance(n, (int, float)) and n > 0]
    return round(sum(nums) / len(nums), 2) if nums else 0.0


@app.route("/api/admin/rankings")
@admin_required
def api_admin_rankings():
    subs = Submission.query.all()
    group_member_avgs = {g: [] for g in GROUPS}
    group_rubric_scores = {g: {r: [] for r in RUBRIC} for g in GROUPS}
    per_member = []
    for s in subs:
        payload = json.loads(s.payload or "{}")
        groups_data = payload.get("groups", {})
        member_entry = {"code": s.member_code, "name": s.member_name, "groups": {}, "rubric_detail": {}}
        for g in GROUPS:
            gd = groups_data.get(g, {})
            scores = gd.get("q3", {}) or {}
            vals = [scores.get(k, 0) for k in RUBRIC]
            avg = _avg(vals)
            if avg > 0:
                group_member_avgs[g].append(avg)
            member_entry["groups"][g] = avg
            rd = {}
            for r in RUBRIC:
                v = scores.get(r, 0)
                if isinstance(v, (int, float)) and v > 0:
                    group_rubric_scores[g][r].append(v)
                rd[r] = v
            member_entry["rubric_detail"][g] = rd
        per_member.append(member_entry)

    rankings = []
    for g in GROUPS:
        score = _avg(group_member_avgs[g])
        rubric = {}
        for r in RUBRIC:
            a = _avg(group_rubric_scores[g][r])
            rubric[r] = {"avg": a, "pct": round(a / 5 * 100, 1)}
        rankings.append({
            "group": g,
            "score": score,
            "pct": round(score / 5 * 100, 1),
            "voters": len(group_member_avgs[g]),
            "rubric": rubric,
        })
    rankings.sort(key=lambda x: x["score"], reverse=True)

    # 收集 Q5 / Q6 整體評價
    best_votes = {}
    worst_votes = {}
    self_reviews = []
    best_reasons = []
    worst_reasons = []
    for s in subs:
        payload = json.loads(s.payload or "{}")
        overall = payload.get("overall", {})
        b = overall.get("best_group")
        w = overall.get("worst_group")
        if b:
            best_votes[b] = best_votes.get(b, 0) + 1
            if overall.get("best_reason"):
                best_reasons.append({"by": s.member_name, "group": b, "reason": overall.get("best_reason")})
        if w:
            worst_votes[w] = worst_votes.get(w, 0) + 1
            if overall.get("worst_reason"):
                worst_reasons.append({"by": s.member_name, "group": w, "reason": overall.get("worst_reason")})
        if overall.get("self_review"):
            self_reviews.append({"by": s.member_name, "text": overall.get("self_review")})

    return jsonify({
        "rankings": rankings,
        "per_member": per_member,
        "best_votes": best_votes,
        "worst_votes": worst_votes,
        "best_reasons": best_reasons,
        "worst_reasons": worst_reasons,
        "self_reviews": self_reviews,
    })


AI_SYSTEM_PROMPT = """你是課堂同儕評分的整合助理。第2組5位成員針對第1～6組（共6組）進行課堂表現評分。
你的任務是把多位成員對同一題的意見整合成一段簡潔的綜合意見。

請以**純 JSON** 格式回覆，不要有任何 Markdown、code block 或其他文字，schema 必須完全符合：
{
  "groups": {
    "第1組": {"highlights_summary": "一句話總結亮點（15~25字）", "highlights": "詳細綜合亮點（60~140字）", "improvements_summary": "一句話總結優化建議（15~25字）", "improvements": "詳細綜合優化建議（60~140字）"},
    "第2組": {"highlights_summary": "…", "highlights": "…", "improvements_summary": "…", "improvements": "…"},
    "第3組": {"highlights_summary": "…", "highlights": "…", "improvements_summary": "…", "improvements": "…"},
    "第4組": {"highlights_summary": "…", "highlights": "…", "improvements_summary": "…", "improvements": "…"},
    "第5組": {"highlights_summary": "…", "highlights": "…", "improvements_summary": "…", "improvements": "…"},
    "第6組": {"highlights_summary": "…", "highlights": "…", "improvements_summary": "…", "improvements": "…"}
  },
  "best_overall": "…",
  "worst_overall": "…",
  "self_reflection": "…"
}

寫作要求：
- 用繁體中文
- highlights_summary / improvements_summary 必須是一句話精準總結（15~25字），抓出最核心的重點
- highlights / improvements 是詳細版本（60~140字），找出多位成員提到的共同點，用「多數成員認為…」「部分成員指出…」呈現
- 若意見少於 2 則，誠實寫出單一意見即可
- 若該題完全沒人填寫，寫「（無資料）」
- 客觀中立，避免攻擊性用詞
- 直接給結論，不要說「綜合來看」之類的引言"""


@app.route("/api/admin/ai_summary", methods=["GET"])
@admin_required
def api_get_ai_summary():
    s = AISummary.query.order_by(AISummary.id.desc()).first()
    if not s:
        return jsonify({"summary": None, "created_at": None})
    return jsonify({
        "summary": json.loads(s.payload),
        "created_at": s.created_at.strftime("%Y-%m-%d %H:%M:%S"),
    })


@app.route("/api/admin/ai_summary", methods=["POST"])
@admin_required
def api_generate_ai_summary():
    if not anthropic_client:
        return jsonify({"error": "伺服器未設定 ANTHROPIC_API_KEY"}), 500

    subs = Submission.query.all()
    by_group = {g: {"q1": [], "q2": []} for g in GROUPS}
    best_reasons, worst_reasons, self_reviews = [], [], []
    for s in subs:
        p = json.loads(s.payload or "{}")
        for g in GROUPS:
            gd = (p.get("groups") or {}).get(g, {})
            if (gd.get("q1") or "").strip():
                by_group[g]["q1"].append(f"{s.member_name}：{gd['q1'].strip()}")
            if (gd.get("q2") or "").strip():
                by_group[g]["q2"].append(f"{s.member_name}：{gd['q2'].strip()}")
        ov = p.get("overall") or {}
        if ov.get("best_group") and (ov.get("best_reason") or "").strip():
            best_reasons.append(f"{s.member_name} 推薦 {ov['best_group']}：{ov['best_reason'].strip()}")
        if ov.get("worst_group") and (ov.get("worst_reason") or "").strip():
            worst_reasons.append(f"{s.member_name} 認為 {ov['worst_group']} 需優化：{ov['worst_reason'].strip()}")
        if (ov.get("self_review") or "").strip():
            self_reviews.append(f"{s.member_name}：{ov['self_review'].strip()}")

    sections = []
    for g in GROUPS:
        h = "\n".join(f"- {x}" for x in by_group[g]["q1"]) or "（無）"
        i = "\n".join(f"- {x}" for x in by_group[g]["q2"]) or "（無）"
        sections.append(f"## {g}\n【亮點意見】\n{h}\n【優化意見】\n{i}")
    sections.append("## Q5a 最優秀組別投票理由\n" + ("\n".join(f"- {x}" for x in best_reasons) or "（無）"))
    sections.append("## Q5b 最需優化組別投票理由\n" + ("\n".join(f"- {x}" for x in worst_reasons) or "（無）"))
    sections.append("## Q6 第2組（自組）反思\n" + ("\n".join(f"- {x}" for x in self_reviews) or "（無）"))
    user_content = "\n\n".join(sections)

    try:
        msg = anthropic_client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=4096,
            system=AI_SYSTEM_PROMPT,
            messages=[{"role": "user", "content": user_content}],
        )
    except Exception as e:
        return jsonify({"error": f"Anthropic API 錯誤：{e}"}), 500

    text = msg.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1].rsplit("```", 1)[0].strip()
        if text.startswith("json"):
            text = text[4:].strip()
    try:
        result = json.loads(text)
    except Exception:
        return jsonify({"error": "AI 回傳非 JSON 格式", "raw": text}), 500

    rec = AISummary(payload=json.dumps(result, ensure_ascii=False))
    db.session.add(rec)
    db.session.commit()
    # 只保留最近 10 筆
    old = AISummary.query.order_by(AISummary.id.desc()).offset(10).all()
    for o in old:
        db.session.delete(o)
    db.session.commit()
    return jsonify({"summary": result, "created_at": rec.created_at.strftime("%Y-%m-%d %H:%M:%S")})


@app.route("/api/admin/export_ppt")
@admin_required
def api_admin_export_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    subs = Submission.query.all()
    group_rubric_scores = {g: {r: [] for r in RUBRIC} for g in GROUPS}
    group_member_avgs = {g: [] for g in GROUPS}
    per_member = []
    for s in subs:
        payload = json.loads(s.payload or "{}")
        groups_data = payload.get("groups", {})
        entry = {"name": s.member_name, "groups": {}}
        for g in GROUPS:
            gd = groups_data.get(g, {})
            scores = gd.get("q3", {}) or {}
            vals = [scores.get(k, 0) for k in RUBRIC]
            avg = _avg(vals)
            if avg > 0:
                group_member_avgs[g].append(avg)
            entry["groups"][g] = avg
            for r in RUBRIC:
                v = scores.get(r, 0)
                if isinstance(v, (int, float)) and v > 0:
                    group_rubric_scores[g][r].append(v)
        per_member.append(entry)

    rankings = []
    for g in GROUPS:
        score = _avg(group_member_avgs[g])
        rubric = {}
        for r in RUBRIC:
            a = _avg(group_rubric_scores[g][r])
            rubric[r] = a
        rankings.append({"group": g, "score": score, "pct": round(score / 5 * 100, 1), "rubric": rubric})
    rankings.sort(key=lambda x: x["score"], reverse=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BG = RGBColor(0xF5, 0xF7, 0xFB)
    PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    MUTED = RGBColor(0x6B, 0x72, 0x80)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    WARN = RGBColor(0xF5, 0x9E, 0x0B)
    RED = RGBColor(0xDC, 0x26, 0x26)

    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    def color_for_pct(pct):
        if pct >= 80:
            return GREEN
        elif pct >= 60:
            return WARN
        return RED

    def set_cell(cell, text, size=11, bold=False, color=DARK, align=PP_ALIGN.CENTER):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        cell.margin_left = Emu(36000)
        cell.margin_right = Emu(36000)
        cell.margin_top = Emu(18000)
        cell.margin_bottom = Emu(18000)

    def fill_cell(cell, color):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    # === Slide 1: Title ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide1, BG)
    add_text(slide1, 1, 2.2, 11, 1.2, "第2組 · 同儕評分結果", size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide1, 1, 3.6, 11, 0.6, "課堂分組報告互評 · 結構化五維度評分", size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide1, 1, 4.5, 11, 0.5, datetime.utcnow().strftime("%Y-%m-%d"), size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # === Slide 2: Rankings + Rubric ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, BG)
    add_text(slide2, 0.5, 0.3, 12, 0.5, "各組綜合評分排名 · 五維度分析", size=24, bold=True, color=PRIMARY)

    cols = ["排名", "組別"] + RUBRIC + ["平均", "百分比"]
    n_cols = len(cols)
    n_rows = len(rankings) + 1
    col_widths = [0.6, 1.0] + [1.3] * 5 + [1.0, 1.0]
    total_w = sum(col_widths)
    tbl_left = (13.333 - total_w) / 2

    table_shape = slide2.shapes.add_table(n_rows, n_cols, Inches(tbl_left), Inches(1.0), Inches(total_w), Inches(0.45 * n_rows))
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    for ci, h in enumerate(cols):
        cell = table.cell(0, ci)
        set_cell(cell, h, size=10, bold=True, color=WHITE)
        fill_cell(cell, PRIMARY)

    for ri, rk in enumerate(rankings):
        row = ri + 1
        set_cell(table.cell(row, 0), str(ri + 1), size=14, bold=True, color=PRIMARY)
        set_cell(table.cell(row, 1), rk["group"], size=12, bold=True)
        for di, dim in enumerate(RUBRIC):
            a = rk["rubric"][dim]
            pct = round(a / 5 * 100, 1)
            set_cell(table.cell(row, 2 + di), f"{a:.2f}\n({pct:.0f}%)", size=10, color=color_for_pct(pct))
        set_cell(table.cell(row, 7), f"{rk['score']:.2f}", size=12, bold=True, color=PRIMARY)
        set_cell(table.cell(row, 8), f"{rk['pct']:.1f}%", size=12, bold=True, color=color_for_pct(rk["pct"]))
        if ri % 2 == 0:
            for ci in range(n_cols):
                fill_cell(table.cell(row, ci), RGBColor(0xEE, 0xF2, 0xFF))

    # === Slide 3: Per-member scores ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, BG)
    add_text(slide3, 0.5, 0.3, 12, 0.5, "個人評分一覽（每位成員對各組的平均分）", size=24, bold=True, color=PRIMARY)

    sorted_groups = [r["group"] for r in rankings]
    m_cols = ["成員"] + sorted_groups
    n_mc = len(m_cols)
    n_mr = len(per_member) + 1
    mc_widths = [1.6] + [1.5] * len(sorted_groups)
    mc_total = sum(mc_widths)
    mc_left = (13.333 - mc_total) / 2

    m_shape = slide3.shapes.add_table(n_mr, n_mc, Inches(mc_left), Inches(1.0), Inches(mc_total), Inches(0.5 * n_mr))
    m_table = m_shape.table
    for ci, w in enumerate(mc_widths):
        m_table.columns[ci].width = Inches(w)

    for ci, h in enumerate(m_cols):
        cell = m_table.cell(0, ci)
        set_cell(cell, h, size=11, bold=True, color=WHITE)
        fill_cell(cell, PRIMARY)

    for mi, m in enumerate(per_member):
        row = mi + 1
        set_cell(m_table.cell(row, 0), m["name"], size=12, bold=True)
        for gi, g in enumerate(sorted_groups):
            avg = m["groups"].get(g, 0)
            pct = round(avg / 5 * 100, 1) if avg else 0
            set_cell(m_table.cell(row, 1 + gi), f"{avg:.2f}\n({pct:.0f}%)", size=11, color=color_for_pct(pct) if avg > 0 else MUTED)
        if mi % 2 == 0:
            for ci in range(n_mc):
                fill_cell(m_table.cell(row, ci), RGBColor(0xEE, 0xF2, 0xFF))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    fname = f"tempform_ranking_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
    return send_file(buf, as_attachment=True, download_name=fname, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.route("/healthz")
def healthz():
    return {"ok": True}


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5205)), debug=True)
